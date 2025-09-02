# tests/create_test_files.py
import os
from pathlib import Path
import sqlite3

out = Path("tests/fixtures")
out.mkdir(parents=True, exist_ok=True)

# 1) sample.txt
(out / "sample.txt").write_text("This is a sample text file.\nIt contains two lines.\n", encoding="utf-8")

# 2) sample.md
(out / "sample.md").write_text("# Sample MD\n\nThis is markdown content.\n\n|A|B|\n|--:|--|\n|1|x|\n|2|y|\n", encoding="utf-8")

# 3) sample.csv
(out / "sample.csv").write_text("name,age\nAlice,30\nBob,28\n", encoding="utf-8")

# 4) sample.xlsx (pandas + openpyxl required)
try:
    import pandas as pd
    df = pd.DataFrame({"Product":["Laptop","Phone"], "Sales":[1200,800]})
    df.to_excel(out / "sample.xlsx", index=False, sheet_name="Sheet1")
except Exception as e:
    print("Skipping xlsx creation:", e)

# 5) sample.docx (python-docx)
try:
    from docx import Document
    doc = Document()
    doc.add_heading('Sample DOCX', level=1)
    doc.add_paragraph('This is a paragraph in the docx file.')
    t = doc.add_table(rows=1, cols=2)
    hdr = t.rows[0].cells
    hdr[0].text = "ColA"; hdr[1].text = "ColB"
    row = t.add_row().cells
    row[0].text = "1"; row[1].text = "x"
    doc.save(out / "sample.docx")
except Exception as e:
    print("Skipping docx creation:", e)

# 6) sample.pptx (python-pptx)
try:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(left=1000000, top=1000000, width=6000000, height=2000000)
    tf = tx.text_frame
    tf.text = "This is a sample pptx slide text"
    prs.save(out / "sample.pptx")
except Exception as e:
    print("Skipping pptx creation:", e)

# 7) sample.pdf (PyMuPDF / fitz)
try:
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Sample PDF created for tests.\nLine 2.")
    doc.save(str(out / "sample.pdf"))
    doc.close()
except Exception as e:
    print("Skipping pdf creation:", e)

# 8) sample.sqlite (sqlite3)
try:
    dbp = out / "sample.sqlite"
    conn = sqlite3.connect(str(dbp))
    cur = conn.cursor()
    cur.execute("CREATE TABLE IF NOT EXISTS people (id INTEGER PRIMARY KEY, name TEXT, age INTEGER)")
    cur.execute("DELETE FROM people")
    cur.executemany("INSERT INTO people (name, age) VALUES (?, ?)", [("Alice", 30), ("Bob", 28)])
    conn.commit()
    conn.close()
except Exception as e:
    print("Skipping sqlite creation:", e)

print("Test fixtures creation attempted. Check tests/fixtures/")
