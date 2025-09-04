import os
import sys
import pandas as pd
from typing import List
from langchain.schema import Document
from langchain_core.messages import HumanMessage
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import camelot
from zipfile import ZipFile
from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException
import sqlite3

log = CustomLogger().get_logger(__name__)

def extract_tables(file_path: str, extension: str) -> List[Document]:
    """
    Extracts tables from the document and converts them to markdown for inclusion as Documents.
    """
    tables: List[Document] = []
    try:
        if extension == '.pdf':
            table_list = camelot.read_pdf(str(file_path), flavor='stream', pages='all')
            for idx, table in enumerate(table_list):
                df = table.df
                content = f"Extracted Table {idx + 1}:\n{df.to_markdown()}"
                tables.append(Document(page_content=content, metadata={"source": file_path, "type": "table", "index": idx}))

        elif extension == '.docx':
            doc = DocxDocument(file_path)
            for idx, table in enumerate(doc.tables):
                data = [[cell.text for cell in row.cells] for row in table.rows]
                if data:
                    df = pd.DataFrame(data[1:], columns=data[0])
                    content = f"Extracted Table {idx + 1}:\n{df.to_markdown()}"
                    tables.append(Document(page_content=content, metadata={"source": file_path, "type": "table", "index": idx}))

        elif extension == '.pptx':
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_table:
                        table = shape.table
                        data = [[cell.text for cell in row.cells] for row in table.rows]
                        if data:
                            df = pd.DataFrame(data[1:], columns=data[0])
                            content = f"Extracted Table in slide {slide_idx + 1}:\n{df.to_markdown()}"
                            tables.append(Document(page_content=content, metadata={"source": file_path, "type": "table", "slide": slide_idx}))

        elif extension == '.xlsx':
            df = pd.read_excel(file_path)
            content = f"Excel Table:\n{df.to_markdown()}"
            tables.append(Document(page_content=content, metadata={"source": file_path, "type": "table"}))

        elif extension == '.csv':
            df = pd.read_csv(file_path)
            content = f"CSV Table:\n{df.to_markdown()}"
            tables.append(Document(page_content=content, metadata={"source": file_path, "type": "table"}))

        log.info("Tables extracted successfully", file_path=file_path, count=len(tables))
        return tables

    except Exception as e:
        log.error(f"Table extraction failed for {file_path}", error=str(e))
        raise DocumentPortalException("Table extraction failed", sys) from e

def extract_and_describe_images(file_path: str, extension: str, llm, provider: str) -> List[Document]:
    """
    Extracts images from the document, describes them using the LLM (if provider supports multimodal), and creates Documents.
    """
    descriptions: List[Document] = []
    if provider != 'google':
        log.info("Image description skipped: Provider does not support multimodal")
        return descriptions

    temp_dir = 'temp_images'
    os.makedirs(temp_dir, exist_ok=True)
    try:
        if extension == '.pdf':
            doc = fitz.open(file_path)
            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    if base_image:
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        image_path = os.path.join(temp_dir, f"pdf_page{page_num}_img{img_index}.{image_ext}")
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                        message = HumanMessage(content=[
                            {"type": "text", "text": "Describe this image in detail for document analysis."},
                            {"type": "image_url", "image_url": image_path}
                        ])
                        response = llm.invoke([message])
                        desc = response.content
                        descriptions.append(Document(page_content=f"Image description on page {page_num+1}: {desc}", metadata={"source": file_path, "type": "image", "page": page_num}))
                        os.remove(image_path)
            doc.close()

        elif extension == '.docx':
            with ZipFile(file_path, 'r') as zip_ref:
                for file in zip_ref.namelist():
                    if file.startswith('word/media/'):
                        image_bytes = zip_ref.read(file)
                        image_path = os.path.join(temp_dir, os.path.basename(file))
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                        message = HumanMessage(content=[
                            {"type": "text", "text": "Describe this image in detail for document analysis."},
                            {"type": "image_url", "image_url": image_path}
                        ])
                        response = llm.invoke([message])
                        desc = response.content
                        descriptions.append(Document(page_content=f"Image description: {desc}", metadata={"source": file_path, "type": "image"}))
                        os.remove(image_path)

        elif extension == '.pptx':
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        image_path = os.path.join(temp_dir, f"ppt_slide{slide_idx}_shape{shape_idx}.{image_ext}")
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                        message = HumanMessage(content=[
                            {"type": "text", "text": "Describe this image in detail for document analysis."},
                            {"type": "image_url", "image_url": image_path}
                        ])
                        response = llm.invoke([message])
                        desc = response.content
                        descriptions.append(Document(page_content=f"Image description in slide {slide_idx + 1}: {desc}", metadata={"source": file_path, "type": "image", "slide": slide_idx}))
                        os.remove(image_path)

        log.info("Images extracted and described successfully", file_path=file_path, count=len(descriptions))
        return descriptions

    except Exception as e:
        log.error(f"Image extraction and description failed for {file_path}", error=str(e))
        raise DocumentPortalException("Image extraction failed", sys) from e
    finally:
        if os.path.exists(temp_dir) and not os.listdir(temp_dir):
            os.rmdir(temp_dir)

def load_sqlite_db(file_path: str) -> List[Document]:
    """
    Loads tables from a SQLite .db file as Documents with schema and sample data.
    """
    try:
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        docs: List[Document] = []
        for table in tables:
            table_name = table[0]
            df = pd.read_sql_query(f"SELECT * FROM {table_name} LIMIT 100", conn)  # Limit to prevent large data
            content = f"Database Table: {table_name}\nColumns: {', '.join(df.columns)}\nSample Data:\n{df.to_markdown()}"
            docs.append(Document(page_content=content, metadata={"source": file_path, "type": "db_table", "table": table_name}))
        conn.close()
        log.info("SQLite DB loaded successfully", file_path=file_path, table_count=len(docs))
        return docs
    except Exception as e:
        log.error(f"SQLite DB load failed for {file_path}", error=str(e))
        raise DocumentPortalException("SQLite DB load failed", sys) from e