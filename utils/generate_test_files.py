import os
import pandas as pd
import sqlite3
from pptx import Presentation
from pptx.util import Inches
from PIL import Image  # For generating a sample image
from io import BytesIO
import nltk

# Download NLTK punkt tokenizer to avoid runtime download errors
try:
    nltk.download('punkt', quiet=True)
    print("NLTK punkt tokenizer downloaded successfully")
except Exception as e:
    print(f"Failed to download NLTK punkt tokenizer: {e}")

# Create the data/test_files folder if it doesn't exist
os.makedirs('data/test_files', exist_ok=True)

# 1. Generate .md (Markdown file with text)
md_content = """
# Sample Markdown Document

This is a test Markdown file.

- Item 1
- Item 2

**Bold text** and *italic text*.
"""
with open('data/test_files/test.md', 'w') as f:
    f.write(md_content)

# 2. Generate .xlsx (Excel file with table)
df = pd.DataFrame({
    "Name": ["Alice", "Bob", "Charlie"],
    "Age": [25, 30, 35],
    "City": ["New York", "London", "Paris"]
})
df.to_excel('data/test_files/test.xlsx', index=False)

# 3. Generate .csv (CSV file with table)
df.to_csv('data/test_files/test.csv', index=False)

# 4. Generate .sql (SQL script file)
sql_content = """
-- Sample SQL Script

CREATE TABLE users (
    id INTEGER PRIMARY KEY,
    name TEXT,
    age INTEGER,
    city TEXT
);

INSERT INTO users (name, age, city) VALUES ('Alice', 25, 'New York');
INSERT INTO users (name, age, city) VALUES ('Bob', 30, 'London');
INSERT INTO users (name, age, city) VALUES ('Charlie', 35, 'Paris');

SELECT * FROM users;
"""
with open('data/test_files/test.sql', 'w') as f:
    f.write(sql_content)

# 5. Generate .db (SQLite database with table)
conn = sqlite3.connect('data/test_files/test.db')
cursor = conn.cursor()
cursor.execute("""
CREATE TABLE users (
    id INTEGER PRIMARY KEY,
    name TEXT,
    age INTEGER,
    city TEXT
)
""")
cursor.execute("INSERT INTO users (name, age, city) VALUES ('Alice', 25, 'New York')")
cursor.execute("INSERT INTO users (name, age, city) VALUES ('Bob', 30, 'London')")
cursor.execute("INSERT INTO users (name, age, city) VALUES ('Charlie', 35, 'Paris')")
conn.commit()
conn.close()

# 6. Generate .pptx (PowerPoint with text, table, and image)
prs = Presentation()

# Slide 1: Text
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
title.text = "Sample PowerPoint"
subtitle = slide1.placeholders[1]
subtitle.text = "This is a test slide with text."

# Slide 2: Table
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
title2 = slide2.shapes.title
title2.text = "Slide with Table"
rows, cols = 4, 3
left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(4)
table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
table.cell(0, 0).text = "Name"
table.cell(0, 1).text = "Age"
table.cell(0, 2).text = "City"
table.cell(1, 0).text = "Alice"
table.cell(1, 1).text = "25"
table.cell(1, 2).text = "New York"
table.cell(2, 0).text = "Bob"
table.cell(2, 1).text = "30"
table.cell(2, 2).text = "London"
table.cell(3, 0).text = "Charlie"
table.cell(3, 1).text = "35"
table.cell(3, 2).text = "Paris"

# Slide 3: Image
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
title3 = slide3.shapes.title
title3.text = "Slide with Image"

# Generate a sample image with Pillow (simple red square)
img = Image.new('RGB', (200, 200), color='red')
img_buffer = BytesIO()
img.save(img_buffer, format='PNG')
img_buffer.seek(0)

# Add the image to the slide
left_img = Inches(2)
top_img = Inches(2)
slide3.shapes.add_picture(img_buffer, left_img, top_img, width=Inches(4))

prs.save('data/test_files/test.pptx')

# Print generated files for verification
generated_files = os.listdir('data/test_files')
print("Generated test files:", generated_files)