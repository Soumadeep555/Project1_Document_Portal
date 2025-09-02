# src/document_ingestion/loaders/ppt_loader.py
from pathlib import Path
from pptx import Presentation
from .base import BaseLoader, LoadedDocument, ImageBlob
from .image_utils import save_image_bytes


class PptLoader(BaseLoader):
    def load(self, path: str) -> LoadedDocument:
        p = Path(path)
        prs = Presentation(str(p))
        texts = []
        images = []
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                # images: shape may have image attribute via shape.image
                try:
                    img = getattr(shape, "image", None)
                    if img is not None and img.blob:
                        fname = getattr(img, "filename", f"{p.stem}_s{si}.png")
                        out = save_image_bytes(img.blob, f"{p.stem}_s{si}_{fname}")
                        images.append(ImageBlob(path=out, source=str(p), page_or_index=si))
                except Exception:
                    pass
        return LoadedDocument(source=str(p), text="\n".join(texts), images=images)
