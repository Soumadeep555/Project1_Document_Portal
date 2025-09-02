# src/document_ingestion/data_ingestion.py
import os
import tempfile
from typing import List, Any, Optional
from pathlib import Path
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import markdown
import fitz  # PyMuPDF
from PIL import Image
import io
import csv
import sqlite3
import docx2txt

# -------------------------------
# DocHandler class
# -------------------------------
class DocHandler:
    """
    Unified document handler to read text, tables, and images from various document types.
    """

    def __init__(self, path: str):
        self.path = path
        self.extension = Path(path).suffix.lower()
        self.text: str = ""
        self.tables: List[Any] = []
        self.images: List[Image.Image] = []

        self._load_document()

    # -------------------------------
    # Private loader
    # -------------------------------
    def _load_document(self):
        if self.extension == ".txt":
            self._load_txt()
        elif self.extension == ".pdf":
            self._load_pdf()
        elif self.extension == ".docx":
            self._load_docx()
        elif self.extension == ".pptx":
            self._load_pptx()
        elif self.extension == ".md":
            self._load_md()
        elif self.extension in [".xlsx", ".xls"]:
            self._load_excel()
        elif self.extension == ".csv":
            self._load_csv()
        elif self.extension in [".db", ".sqlite", ".sqlite3"]:
            self._load_sql()
        else:
            raise ValueError(f"Unsupported file type: {self.extension}")

    # -------------------------------
    # Individual loaders
    # -------------------------------
    def _load_txt(self):
        with open(self.path, "r", encoding="utf-8") as f:
            self.text = f.read()

    def _load_pdf(self):
        doc = fitz.open(self.path)
        text_chunks = []
        for page in doc:
            text_chunks.append(page.get_text())
            # extract images
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes))
                self.images.append(image)
        self.text = "\n".join(text_chunks)

    def _load_docx(self):
        self.text = docx2txt.process(self.path)
        doc = DocxDocument(self.path)
        # extract tables
        for table in doc.tables:
            table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            self.tables.append(table_data)
        # extract images
        for rel in doc.part._rels:
            rel_obj = doc.part._rels[rel]
            if "image" in rel_obj.target_ref:
                image_bytes = rel_obj.target_part.blob
                image = Image.open(io.BytesIO(image_bytes))
                self.images.append(image)

    def _load_pptx(self):
        prs = Presentation(self.path)
        text_chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_chunks.append(shape.text)
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        self.images.append(Image.open(io.BytesIO(image_bytes)))
                    except Exception:
                        pass
        self.text = "\n".join(text_chunks)

    def _load_md(self):
        with open(self.path, "r", encoding="utf-8") as f:
            md_text = f.read()
        # Convert Markdown to plain text
        self.text = markdown.markdown(md_text)

    def _load_excel(self):
        xls = pd.ExcelFile(self.path)
        self.tables = [xls.parse(sheet_name).fillna("").values.tolist() for sheet_name in xls.sheet_names]
        self.text = "\n".join([str(table) for table in self.tables])

    def _load_csv(self):
        df = pd.read_csv(self.path)
        self.tables.append(df.fillna("").values.tolist())
        self.text = df.to_csv(index=False)

    def _load_sql(self):
        conn = sqlite3.connect(self.path)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        for table_name in tables:
            df = pd.read_sql_query(f"SELECT * FROM {table_name[0]}", conn)
            self.tables.append(df.fillna("").values.tolist())
            self.text += df.to_csv(index=False) + "\n"
        conn.close()

    # -------------------------------
    # Optional save helper
    # -------------------------------
    @staticmethod
    def save_file(uploaded_file, save_dir: Optional[str] = None) -> str:
        save_dir = save_dir or tempfile.gettempdir()
        os.makedirs(save_dir, exist_ok=True)
        path = os.path.join(save_dir, uploaded_file.filename)
        with open(path, "wb") as f:
            f.write(uploaded_file.file.read())
        return path
